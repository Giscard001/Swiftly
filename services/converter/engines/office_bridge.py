"""Conversions docx <-> pptx via python-docx + python-pptx.

LibreOffice ne sait pas convertir directement docx<->pptx (le bridge via ODP
échoue en silence). On utilise ici les bibliothèques Python dédiées :
  - docx -> pptx : lire les paragraphes du docx, créer un slide par paragraphe titre
  - pptx -> docx : lire le texte des slides, créer un paragraphe docx par slide

La fidélité est limitée (Word et PowerPoint n'ont pas la même structure),
mais la conversion est robuste et préserve le texte.
"""
from .base import ConversionRoute, ConvertContext

try:
    from docx import Document
    _DOCX = True
except Exception:
    _DOCX = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    _PPTX = True
except Exception:
    _PPTX = False

NAME = "office_bridge"


def available() -> bool:
    return _DOCX and _PPTX


def routes() -> list[ConversionRoute]:
    if not available():
        return []
    return [
        ConversionRoute("document", "docx", "pptx", NAME, _docx_to_pptx, "DOCX -> PPTX"),
        ConversionRoute("document", "pptx", "docx", NAME, _pptx_to_docx, "PPTX -> DOCX"),
    ]


def _docx_to_pptx(ctx: ConvertContext) -> None:
    """Word -> PowerPoint : un slide par titre de section + ses paragraphes.
    Construction 100% contrôlée (layout vierge + textboxes) pour éviter le
    chevauchement de texte des placeholders de template par défaut."""
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

    ctx.set_progress(10, "Lecture du document Word")
    doc = Document(str(ctx.input_path))

    # Récupère les paragraphes avec leur style
    paragraphs = []
    for p in doc.paragraphs:
        text = (p.text or "").strip()
        if not text:
            continue
        is_heading = p.style is not None and "Heading" in (p.style.name or "")
        paragraphs.append((text, is_heading))
    if not paragraphs:
        raise RuntimeError("Document Word vide (aucun paragraphe de texte)")

    ctx.set_progress(30, "Construction des slides")
    prs = Presentation()
    # Force 16:9 pour avoir de la place (13.33 x 7.5 pouces par défaut en EMU)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    # Layout vierge (index 6 dans le template par défaut python-pptx)
    blank_layout = prs.slide_layouts[6]

    # Groupe les paragraphes : un nouveau slide démarre à chaque titre
    slides_content: list[tuple[str, list[str]]] = []
    current_title = ""
    current_body: list[str] = []

    for text, is_heading in paragraphs:
        if is_heading or (not current_title and not current_body):
            if current_title or current_body:
                slides_content.append((current_title, current_body))
            current_title = text
            current_body = []
        else:
            current_body.append(text)
            # découpe les longues sections sans titre (limite par slide)
            if len(current_body) >= 6:
                slides_content.append((current_title, current_body))
                current_title = ""
                current_body = []
    if current_title or current_body:
        slides_content.append((current_title, current_body))

    if not slides_content:
        slides_content = [("", [p[0] for p in paragraphs[:6]])]

    total = len(slides_content)
    for i, (title, body) in enumerate(slides_content):
        slide = prs.slides.add_slide(blank_layout)

        # ── Zone de titre (en haut, pleine largeur) ───────────────────────
        if title:
            tbox = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4), SW - Inches(1.0), Inches(1.1)
            )
            tf = tbox.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE  # taille fixe, on wrap le texte
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = title[:200]
            run.font.size = Pt(32)
            run.font.bold = True

        # ── Zone de corps (sous le titre, pleine largeur) ─────────────────
        if body:
            bbox = slide.shapes.add_textbox(
                Inches(0.6),
                Inches(1.8) if title else Inches(0.6),
                SW - Inches(1.2),
                SH - (Inches(2.2) if title else Inches(1.0)),
            )
            tf = bbox.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            # Chaque paragraphe du docx = un paragraphe PPT séparé
            for j, line in enumerate(body[:20]):  # plafond anti-débordement
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(6)
                run = p.add_run()
                run.text = line[:400]
                run.font.size = Pt(18)

        ctx.set_progress(int(30 + 65 * (i + 1) / max(total, 1)), f"Slide {i + 1}/{total}")

    ctx.set_progress(98, "Enregistrement")
    prs.save(str(ctx.output_path))
    ctx.set_progress(100, "Termine")


def _pptx_to_docx(ctx: ConvertContext) -> None:
    """PowerPoint -> Word : un titre de niveau 1 par slide + le texte en paragraphes."""
    ctx.set_progress(10, "Lecture de la présentation")
    prs = Presentation(str(ctx.input_path))
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    total = len(prs.slides)
    if total == 0:
        raise RuntimeError("Présentation vide (aucun slide)")

    for i, slide in enumerate(prs.slides):
        title = ""
        body_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Heuristique : si c'est le titre du slide (placeholder idx 0)
            is_title = (
                shape.is_placeholder
                and shape.placeholder_format.idx == 0
            )
            if is_title and not title:
                title = text
            else:
                body_lines.append(text)

        heading = title or f"Slide {i + 1}"
        doc.add_heading(heading, level=1)
        for line in body_lines:
            for sub in line.splitlines():
                if sub.strip():
                    doc.add_paragraph(sub.strip())
        if not body_lines:
            doc.add_paragraph("(slide sans contenu textuel)")
        ctx.set_progress(int(10 + 85 * (i + 1) / total), f"Slide {i + 1}/{total}")

    ctx.set_progress(98, "Enregistrement")
    doc.save(str(ctx.output_path))
    ctx.set_progress(100, "Termine")
