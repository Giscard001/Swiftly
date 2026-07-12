"""Conversions PDF -> formats Office (.docx, .xlsx, .pptx).

Le sens inverse (Office -> PDF) est géré par engines/office.py via LibreOffice.
Ici on utilise des bibliothèques Python spécialisées :

- PDF -> Word (.docx)  : pdf2docx (préserve paragraphes, tableaux, images)
- PDF -> Excel (.xlsx) : pdfplumber + openpyxl (extrait les tableaux, un onglet/page)
- PDF -> PowerPoint    : pdf2image (poppler) + python-pptx (1 slide = 1 page en image)

Chaque dépendance est soft-importée : si elle manque, la route n'est simplement
pas enregistrée (cf. available()/routes()).
"""
import shutil
import tempfile

from .base import ConversionRoute, ConvertContext

NAME = "pdf_office"

try:
    from pdf2docx import Converter as _Pdf2Docx
    _PDF2DOCX = True
except Exception:
    _PDF2DOCX = False

try:
    import pdfplumber
    _PDFPLUMBER = True
except Exception:
    _PDFPLUMBER = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    _PPTX = True
except Exception:
    _PPTX = False

# pdf2image + poppler (déjà utilisé par engines/pdf_images.py)
try:
    from pdf2image import convert_from_path
    _PDF2IMG = True
except Exception:
    _PDF2IMG = False


def _pdftoppm() -> str | None:
    return shutil.which("pdftoppm")


def available() -> bool:
    # L'engine est dispo si au moins une conversion est faisable.
    return _PDF2DOCX or _PDFPLUMBER or (_PPTX and _PDF2IMG and bool(_pdftoppm()))


def routes() -> list[ConversionRoute]:
    out: list[ConversionRoute] = []
    if _PDF2DOCX:
        out.append(ConversionRoute("document", "pdf", "docx", NAME, _pdf_to_docx, "PDF -> Word"))
    if _PDFPLUMBER:
        out.append(ConversionRoute("document", "pdf", "xlsx", NAME, _pdf_to_xlsx, "PDF -> Excel"))
    if _PPTX and _PDF2IMG and _pdftoppm():
        out.append(ConversionRoute("document", "pdf", "pptx", NAME, _pdf_to_pptx, "PDF -> PowerPoint"))
    return out


def _pdf_to_docx(ctx: ConvertContext) -> None:
    ctx.set_progress(10, "Conversion PDF -> Word")
    cv = _Pdf2Docx(str(ctx.input_path))
    try:
        cv.convert(str(ctx.output_path), start=0, end=None)
    finally:
        cv.close()
    ctx.set_progress(100, "Termine")


def _pdf_to_xlsx(ctx: ConvertContext) -> None:
    import openpyxl

    ctx.set_progress(10, "Extraction des tableaux")
    wb = openpyxl.Workbook()
    # supprime la feuille par défaut vide
    default = wb.active
    wb.remove(default)

    with pdfplumber.open(str(ctx.input_path)) as pdf:
        n = len(pdf.pages)
        has_any_table = False
        for i, page in enumerate(pdf.pages):
            tables = page.extract_tables() or []
            sheet = wb.create_sheet(title=f"Page {i + 1}"[:31])
            if tables:
                has_any_table = True
                row_idx = 1
                for t in tables:
                    for row in t:
                        for col_idx, cell in enumerate(row, start=1):
                            sheet.cell(row=row_idx, column=col_idx, value=(cell if cell is not None else ""))
                        row_idx += 1
                    row_idx += 1  # ligne séparatrice entre tables
            else:
                # fallback : texte brut de la page dans la première colonne
                text = page.extract_text() or ""
                for r, line in enumerate(text.splitlines(), start=1):
                    sheet.cell(row=r, column=1, value=line)
            ctx.set_progress(int(10 + 85 * (i + 1) / max(n, 1)), f"Page {i + 1}/{n}")

        if not has_any_table:
            # on l'indique dans le message final (pas d'échec : on a quand même le texte)
            ctx.set_progress(95, "Aucun tableau detecte - texte exporte a la place")

    wb.save(str(ctx.output_path))
    ctx.set_progress(100, "Termine")


def _pdf_to_pptx(ctx: ConvertContext) -> None:
    """PDF -> PPTX : rend chaque page en image, insère une image par slide.
    Résultat : slides non éditables (images), mais fidèles à la mise en page."""
    ctx.set_progress(10, "Rendu des pages PDF en images")
    tmp = tempfile.mkdtemp()
    try:
        images = convert_from_path(str(ctx.input_path), dpi=150, output_folder=tmp, fmt="png")
        n = len(images)
        ctx.set_progress(50, f"{n} page(s) rendue(s)")

        prs = Presentation()
        # 16:9 par défaut (13.33 x 7.5 pouces)
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        blank_layout = prs.slide_layouts[6]  # layout vierge

        for i, img in enumerate(images):
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(img.filename, 0, 0, width=slide_w, height=slide_h)
            ctx.set_progress(int(50 + 45 * (i + 1) / max(n, 1)), f"Slide {i + 1}/{n}")

        prs.save(str(ctx.output_path))
    finally:
        import shutil as _sh
        _sh.rmtree(tmp, ignore_errors=True)
    ctx.set_progress(100, "Termine")
